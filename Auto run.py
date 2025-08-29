import fitz
from tabulate import tabulate
from pptx import Presentation
import subprocess
import os

def file_to_markdown(input_path, output_path):
    """
    Converts a .docx, .pdf, or .pptx file to a simple Markdown file.
    It determines the file type based on the extension and calls the
    appropriate helper function.

    Args:
        input_path (str): The path to the input file.
        output_path (str): The path to save the output Markdown file.
    """
    # Check if the input file exists
    if not os.path.exists(input_path):
        print(f"Error: The file '{input_path}' does not exist.")
        return

    # Determine the file type by its extension
    file_extension = os.path.splitext(input_path)[1].lower()

    if file_extension == '.docx':
        _docx_to_markdown(input_path, output_path)
    elif file_extension == '.pdf':
        _pdf_to_markdown(input_path, output_path)
    elif file_extension == '.pptx':
        _pptx_to_markdown(input_path, output_path)
    else:
        print(f"Error: Unsupported file type '{file_extension}'. Only .docx, .pdf, and .pptx are supported.")
        return

    print(f"Successfully converted '{input_path}' to '{output_path}'.")

def _docx_to_markdown(docx_path, out_path):
    """
    Converts a DOCX file to a simple Markdown file using Pandoc.
    This is a helper function for `file_to_markdown`.
    """
    try:
        command = ['pandoc', '-s', docx_path, '-o', out_path]
        subprocess.run(command, check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        print(f"An error occurred during Pandoc conversion for DOCX: {e.stderr}")
    except FileNotFoundError:
        print("Error: Pandoc is not installed or not in your system's PATH. Please install it to convert .docx files.")

def _pdf_to_markdown(pdf_path, output_path):
    """
    Converts a PDF file to a simple Markdown file.
    This is a helper function for `file_to_markdown`.
    """
    try:
        doc = fitz.open(pdf_path)
        markdown_content = []

        for page_num, page in enumerate(doc):
            # Extract plain text
            text = page.get_text()
            markdown_content.append(text)

            # Attempt to extract tables
            tables = page.find_tables()
            for table in tables:
                df = table.to_pandas()
                if not df.empty:
                    # Convert the pandas DataFrame to a Markdown table
                    markdown_table = tabulate(df, headers='keys', tablefmt='pipe')
                    markdown_content.append("\n" + markdown_table + "\n")

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(markdown_content))
    except Exception as e:
        print(f"An error occurred during PDF conversion: {e}")

def _pptx_to_markdown(pptx_path, output_path):
    """
    Converts a PPTX file to a simple Markdown file.
    This is a helper function for `file_to_markdown`.
    """
    try:
        prs = Presentation(pptx_path)
        markdown_content = []

        for slide_num, slide in enumerate(prs.slides):
            # Start each slide with a new Markdown section
            markdown_content.append(f"## Slide {slide_num + 1}")

            # Iterate through shapes on the slide
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        # Check for lists (bullet points)
                        if paragraph.level > 0:
                            indentation = "  " * (paragraph.level - 1)
                            markdown_content.append(f"{indentation}* {paragraph.text.strip()}")
                        else:
                            text = paragraph.text.strip()
                            if text:
                                # Heuristic: if text is likely a title, treat as a heading
                                is_title = False
                                for run in paragraph.runs:
                                    if run.font.bold:
                                        is_title = True
                                        break
                                if is_title:
                                    markdown_content.append(f"### {text}")
                                else:
                                    markdown_content.append(text)

        # Join all the content and write to the output file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(markdown_content))
    except Exception as e:
        print(f"An error occurred during PPTX conversion: {e}")

def process_all_files_in_directory(input_dir, output_dir):
    """
    Automatically converts all supported files (.docx, .pdf, .pptx)

    """
    # Create the output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        print(f"Created output directory: {output_dir}")

    # Iterate through all files in the input directory
    for filename in os.listdir(input_dir):
        input_path = os.path.join(input_dir, filename)

        # Skip directories
        if os.path.isdir(input_path):
            continue

        # Create a new output filename with a .md extension
        base_filename = os.path.splitext(filename)[0]
        output_path = os.path.join(output_dir, f"{base_filename}.md")

        # Call the main function to convert the file
        print(f"Attempting to convert: {input_path}")
        file_to_markdown(input_path, output_path)

# Example Usage:
# Note: You need to have 'pre_doc' and 'post_doc' directories created
# and the example files in the 'pre_doc' folder for these to work.
# Uncomment the line below to run the new function.

process_all_files_in_directory('pre_doc', 'post_doc')
